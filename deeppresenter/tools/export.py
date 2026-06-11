"""HTML slide → PPTX conversion using local Chromium + python-pptx."""

import asyncio
import io
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

_PX_TO_EMU = 914400 / 96  # 9525 EMU per pixel (96 dpi)

SLIDE_SIZES = {
    "16:9": (1280, 720),
    "4:3":  (960,  720),
    "A4":   (794,  1123),
}


def _px_emu(px: int) -> int:
    return int(px * _PX_TO_EMU)


async def render_slide_png(
    html_path: str,
    width: int,
    height: int,
    chromium_exe: str,
) -> bytes:
    """Render a single HTML slide to PNG using Chromium headless subprocess."""
    with tempfile.TemporaryDirectory() as tmpdir:
        output_png = Path(tmpdir) / "screenshot.png"
        cmd = [
            chromium_exe,
            "--headless=new",
            "--disable-gpu",
            "--no-sandbox",
            "--disable-dev-shm-usage",
            "--hide-scrollbars",
            f"--window-size={width},{height}",
            f"--screenshot={output_png}",
            f"file://{Path(html_path).resolve()}",
        ]
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.DEVNULL,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            _, stderr = await asyncio.wait_for(proc.communicate(), timeout=30)
        except asyncio.TimeoutError:
            proc.kill()
            raise RuntimeError(f"Chromium timed out rendering {html_path}")

        if not output_png.exists():
            err = stderr.decode("utf-8", errors="replace")[:500]
            raise RuntimeError(f"Chromium did not produce a screenshot. stderr: {err}")

        return output_png.read_bytes()


async def html_slides_to_pptx(
    slides_dir: str,
    output_path: str,
    aspect_ratio: str = "16:9",
) -> str:
    """
    Convert all slide_*.html files in slides_dir to a single PPTX file.
    Returns the path of the created PPTX.
    """
    from deeppresenter.tools.chromium import get_chromium_executable

    chromium_exe = get_chromium_executable()

    slides_path = Path(slides_dir)
    html_files = sorted(slides_path.glob("slide_*.html"))
    if not html_files:
        raise ValueError(f"No slide_*.html files found in {slides_dir}")

    w_px, h_px = SLIDE_SIZES.get(aspect_ratio, SLIDE_SIZES["16:9"])

    prs = Presentation()
    prs.slide_width  = Emu(_px_emu(w_px))
    prs.slide_height = Emu(_px_emu(h_px))
    blank_layout = prs.slide_layouts[6]  # completely blank layout

    for html_file in html_files:
        png_bytes = await render_slide_png(str(html_file), w_px, h_px, chromium_exe)
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(png_bytes),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(output_path)
    return output_path
